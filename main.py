from pptx import Presentation
import os

text = ""
for sunular in os.listdir("sunular"):
    f = open("sunular/"+sunular, "rb")
    prs = Presentation(f)
    print(sunular)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
print(text)
exec(text)